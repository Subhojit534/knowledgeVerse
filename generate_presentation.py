import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_sih_presentation(output_pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # blank layout

    # SIH Template Color Palette
    COLOR_PRIMARY_BLUE  = RGBColor(11, 79, 156)   # #0B4F9C (Footer & Accent Blue)
    COLOR_DARK_TEXT     = RGBColor(26, 29, 32)   # #1A1D20 (Main Text)
    COLOR_MUTED_TEXT    = RGBColor(80, 90, 105)  # #505A69
    COLOR_WHITE         = RGBColor(255, 255, 255)
    COLOR_OVAL_BORDER   = RGBColor(112, 48, 160)  # Purple border matching SIH template oval
    COLOR_TEAL_ACCENT   = RGBColor(0, 168, 232)  # #00A8E8 for Flowchart connectors
    COLOR_CARD_BG       = RGBColor(245, 247, 251) # Subtle grey-blue for flowchart nodes

    def add_common_header_footer(slide, title_text, slide_num):
        # 1. Top Left - Oval Badge "Your Team Name" / M1racles
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5), Inches(0.35), Inches(1.8), Inches(1.0)
        )
        oval.fill.solid()
        oval.fill.fore_color.rgb = COLOR_WHITE
        oval.line.color.rgb = COLOR_OVAL_BORDER
        oval.line.width = Pt(1.5)
        tf_o = oval.text_frame
        tf_o.word_wrap = True
        tf_o.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_o = tf_o.paragraphs[0]
        p_o.text = "Team\nM1racles"
        p_o.font.size = Pt(13)
        p_o.font.bold = True
        p_o.font.color.rgb = COLOR_DARK_TEXT
        p_o.font.name = "Arial"
        p_o.alignment = PP_ALIGN.CENTER

        # 2. Slide Header Title (Center Top)
        title_box = slide.shapes.add_textbox(Inches(2.5), Inches(0.4), Inches(8.333), Inches(0.9))
        tf_title = title_box.text_frame
        tf_title.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_DARK_TEXT
        p_title.font.name = "Times New Roman"
        p_title.alignment = PP_ALIGN.CENTER

        # 3. Top Right - SIH Logo Badge
        logo_box = slide.shapes.add_textbox(Inches(10.8), Inches(0.3), Inches(2.2), Inches(1.0))
        tf_logo = logo_box.text_frame
        tf_logo.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_l1 = tf_logo.paragraphs[0]
        p_l1.text = "SMART INDIA"
        p_l1.font.size = Pt(13)
        p_l1.font.bold = True
        p_l1.font.color.rgb = COLOR_PRIMARY_BLUE
        p_l1.alignment = PP_ALIGN.CENTER
        
        p_l2 = tf_logo.add_paragraph()
        p_l2.text = "HACKATHON 2026"
        p_l2.font.size = Pt(11)
        p_l2.font.bold = True
        p_l2.font.color.rgb = COLOR_DARK_TEXT
        p_l2.alignment = PP_ALIGN.CENTER

        # 4. Bottom Footer Bar (Exact SIH Template match)
        footer_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(6.9), Inches(13.333), Inches(0.6)
        )
        footer_rect.fill.solid()
        footer_rect.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
        footer_rect.line.fill.background()

        tf_footer = footer_rect.text_frame
        tf_footer.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_footer = tf_footer.paragraphs[0]
        p_footer.text = f"@SIH Idea submission- Template                                                                                  {slide_num}"
        p_footer.font.size = Pt(11)
        p_footer.font.color.rgb = COLOR_WHITE
        p_footer.font.name = "Arial"
        p_footer.alignment = PP_ALIGN.CENTER

    def add_bullet(tf, text, level=0, bold=False, size=15, space_before=10, color=COLOR_DARK_TEXT):
        p = tf.add_paragraph() if len(tf.paragraphs[0].text) > 0 else tf.paragraphs[0]
        p.text = text
        p.level = level
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_before = Pt(space_before)
        p.space_after = Pt(4)
        return p

    # ==========================================
    # SLIDE 1: TITLE PAGE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)

    # Top Center Title Box
    h_box = slide1.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.333), Inches(1.3))
    tf_h = h_box.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = "SMART INDIA HACKATHON 2026"
    p_h.font.size = Pt(32)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_PRIMARY_BLUE
    p_h.font.name = "Arial"
    p_h.alignment = PP_ALIGN.CENTER

    p_h2 = tf_h.add_paragraph()
    p_h2.text = "TITLE PAGE"
    p_h2.font.size = Pt(26)
    p_h2.font.bold = True
    p_h2.font.color.rgb = COLOR_DARK_TEXT
    p_h2.font.name = "Times New Roman"
    p_h2.alignment = PP_ALIGN.CENTER

    # Main Details Text Box (Left Aligned Bullet Points)
    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.333), Inches(5.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    add_bullet(tf1, "• Problem Statement ID – SIH2026_EDTECH_042", level=0, bold=True, size=18, space_before=8)
    add_bullet(tf1, "• Problem Statement Title – KnowledgeVerse: Gamified AI-Powered Educational RPG Platform", level=0, bold=True, size=18, space_before=14)
    add_bullet(tf1, "• Theme – EdTech / AI in Education / Immersive Gamified Learning", level=0, bold=True, size=18, space_before=14)
    add_bullet(tf1, "• PS Category – Software", level=0, bold=True, size=18, space_before=14)
    add_bullet(tf1, "• Team ID – M1R_2026_8910", level=0, bold=True, size=18, space_before=14)
    add_bullet(tf1, "• Team Name (Registered on portal) – M1racles", level=0, bold=True, size=18, space_before=14)
    add_bullet(tf1, "   Team Leader: Shreyash Raj  |  Team Member: Subhojit Paul", level=0, bold=False, size=16, space_before=6, color=COLOR_MUTED_TEXT)


    # ==========================================
    # SLIDE 2: IDEA TITLE (PROPOSED SOLUTION)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_common_header_footer(slide2, "IDEA TITLE: KNOWLEDGEVERSE", 2)

    tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    # Main Header
    add_bullet(tf2, "❖ Proposed Solution (Describe your Idea/Solution/Prototype)", level=0, bold=True, size=20, space_before=0, color=COLOR_PRIMARY_BLUE)

    # Bullet 1: Detailed explanation
    add_bullet(tf2, "• Detailed explanation of the proposed solution:", level=0, bold=True, size=17, space_before=14)
    add_bullet(tf2, "  - Interactive 2D RPG World: Transforms academic subjects (Math, CS, Science) into explorable islands and subject academies.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf2, "  - Active Questing Model: Students navigate avatars into subject buildings, interact with AI NPCs, and unlock story-driven learning quests.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf2, "  - Multimodal Learning Integration: Blends 2D visual gameplay, dynamic text lessons, adaptive MCQs, and AI-narrated audio voice explanations.", level=1, bold=False, size=15, space_before=4)

    # Bullet 2: How it addresses the problem
    add_bullet(tf2, "• How it addresses the problem:", level=0, bold=True, size=17, space_before=14)
    add_bullet(tf2, "  - Eliminates Study Fatigue & Boredom: Replaces static textbook reading with active, highly engaging RPG gameplay.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf2, "  - Personalized Learning Velocity: Uses adaptive AI tutoring to adjust quiz difficulty and feedback based on individual student performance.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf2, "  - Sustained Retention & Motivation: Keeps students consistent using RPG leveling, XP rewards, avatar gear shop, and global leaderboards.", level=1, bold=False, size=15, space_before=4)

    # Bullet 3: Innovation and uniqueness
    add_bullet(tf2, "• Innovation and uniqueness of the solution:", level=0, bold=True, size=17, space_before=14)
    add_bullet(tf2, "  - Hybrid Flame 2D + Gemini AI: Seamlessly merges a 60 FPS Flutter game loop with real-time Google Gemini curriculum generation.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf2, "  - ElevenLabs Voice TTS Synthesis: Integrates natural voice narration for dual-coding auditory and visual learning.", level=1, bold=False, size=15, space_before=4)


    # ==========================================
    # SLIDE 3: TECHNICAL APPROACH & USER WORKFLOW
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_common_header_footer(slide3, "TECHNICAL APPROACH", 3)

    # Top Section: Technologies
    tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(2.2))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    add_bullet(tf3, "• Technologies to be used (programming languages, frameworks, hardware):", level=0, bold=True, size=17, space_before=0)
    add_bullet(tf3, "  - Frontend & Game Engine: Flutter SDK (Dart) + Flame 2D Engine (60 FPS game loop, sprite animations, Tiled map rendering, collision detection).", level=1, bold=False, size=14, space_before=3)
    add_bullet(tf3, "  - AI Intelligence Layer: Google Gemini API (Dynamic curriculum generation, prompt engineering, adaptive MCQ creation).", level=1, bold=False, size=14, space_before=3)
    add_bullet(tf3, "  - Voice Synthesis Engine: ElevenLabs API (Streaming real-time text-to-speech voice narration for audio learning).", level=1, bold=False, size=14, space_before=3)
    add_bullet(tf3, "  - Backend & Data Services: Python REST API microservices + SharedPreferences local caching & cloud sync state management.", level=1, bold=False, size=14, space_before=3)

    # Section Header for Workflow
    tb_wf_head = slide3.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.733), Inches(0.4))
    tf_wfh = tb_wf_head.text_frame
    add_bullet(tf_wfh, "• Methodology and process for implementation (User Workflow Flowchart):", level=0, bold=True, size=17, space_before=0)

    # Flowchart Diagram Boxes (Clean, structured horizontal flowchart)
    flow_steps = [
        ("1. App Launch & Auth", "Splash Screen ➔ Login\n➔ Avatar Customization"),
        ("2. Overworld Map", "Explore 2D Archipelago\n➔ Enter Subject Island"),
        ("3. Academy & Quest", "Enter Building ➔ Talk to NPC\n➔ Select Quest Topic"),
        ("4. AI Lesson & Voice", "Gemini AI generates lesson\n➔ ElevenLabs voice TTS"),
        ("5. Adaptive MCQ Quiz", "Solve Quiz ➔ Real-Time Check\n(Pass ➔ XP / Fail ➔ Hint)"),
        ("6. Rewards & Ranking", "Level Up ➔ Unlock Shop Gear\n➔ Sync Global Leaderboard")
    ]

    w_box = Inches(1.75)
    h_box = Inches(2.2)
    gap_x = Inches(0.2)
    start_x = Inches(0.8)
    y_pos = Inches(4.2)

    for idx, (title, desc) in enumerate(flow_steps):
        bx = start_x + idx * (w_box + gap_x)
        
        # Step Card
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y_pos, w_box, h_box)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_PRIMARY_BLUE
        card.line.width = Pt(1.5)

        tb_card = slide3.shapes.add_textbox(bx + Inches(0.05), y_pos + Inches(0.1), w_box - Inches(0.1), h_box - Inches(0.2))
        tf_c = tb_card.text_frame
        tf_c.word_wrap = True
        
        p1 = tf_c.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_PRIMARY_BLUE
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf_c.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_DARK_TEXT
        p2.alignment = PP_ALIGN.CENTER

        # Connector Arrow (except last step)
        if idx < 5:
            arr = slide3.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, 
                bx + w_box + Inches(0.02), y_pos + Inches(0.9), Inches(0.16), Inches(0.35)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_TEAL_ACCENT
            arr.line.fill.background()


    # ==========================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_common_header_footer(slide4, "FEASIBILITY AND VIABILITY", 4)

    tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.5))
    tf4 = tb4.text_frame
    tf4.word_wrap = True

    # Bullet 1: Feasibility
    add_bullet(tf4, "• Analysis of the feasibility of the idea:", level=0, bold=True, size=17, space_before=0)
    add_bullet(tf4, "  - Working Prototype Available: Core codebase fully built in Flutter & Flame, validating 2D gameplay, physics, and AI endpoint integration.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf4, "  - Cloud API Scalability: Offloading heavy LLM and TTS voice synthesis to Gemini and ElevenLabs APIs eliminates local hardware compute demands.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf4, "  - Cross-Platform Compatibility: Single Dart codebase targets Mobile (Android, iOS) and Web browsers without platform re-engineering.", level=1, bold=False, size=15, space_before=4)

    # Bullet 2: Challenges & Risks
    add_bullet(tf4, "• Potential challenges and risks:", level=0, bold=True, size=17, space_before=14)
    add_bullet(tf4, "  - Network Latency: Delays during live LLM text generation and real-time audio synthesis during active gameplay.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf4, "  - Low-Spec Device Performance: Maintaining a smooth 60 FPS game engine loop on budget mobile devices.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf4, "  - Low-Bandwidth Environments: Connectivity bottlenecks for students learning in rural or limited internet areas.", level=1, bold=False, size=15, space_before=4)

    # Bullet 3: Overcoming Strategies
    add_bullet(tf4, "• Strategies for overcoming these challenges:", level=0, bold=True, size=17, space_before=14)
    add_bullet(tf4, "  - Smart Pre-fetching & Caching: Background caching of upcoming quiz modules and audio snippets locally on device.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf4, "  - Game Engine Batching: Flame sprite batching, object pooling, and spatial partitioning for minimal GPU/CPU overhead.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf4, "  - Offline Fallback System: Pre-loaded offline subject question bank ensuring uninterrupted learning when disconnected.", level=1, bold=False, size=15, space_before=4)


    # ==========================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_common_header_footer(slide5, "IMPACT AND BENEFITS", 5)

    tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.5))
    tf5 = tb5.text_frame
    tf5.word_wrap = True

    # Bullet 1: Target Audience Impact
    add_bullet(tf5, "• Potential impact on the target audience:", level=0, bold=True, size=17, space_before=0)
    add_bullet(tf5, "  - 10x Student Engagement: Replaces tedious textbook memorization with rewarding RPG quests and exploration.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf5, "  - Personalized Learning Velocity: Gemini AI tailors question difficulty directly to student response patterns.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf5, "  - Multimodal Retention: Combining 2D visual game graphics, written text, and voice narration maximizes long-term memory recall.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf5, "  - Gamified Mastery: RPG leveling, badges, and shop inventory foster consistent daily study habits.", level=1, bold=False, size=15, space_before=4)

    # Bullet 2: Social & Economic Benefits
    add_bullet(tf5, "• Benefits of the solution (social, economic, environmental, etc.):", level=0, bold=True, size=17, space_before=16)
    add_bullet(tf5, "  - Democratizing 1-on-1 Tutoring: High-quality personalized AI tutoring accessible at zero extra cost to students globally.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf5, "  - Reducing Dropout Rates: Interactive storytelling prevents study boredom and student burnout.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf5, "  - Teacher Assistance Tool: Automated AI progress tracking & analytics relieve educator workload in grading.", level=1, bold=False, size=15, space_before=4)
    add_bullet(tf5, "  - Low Infrastructure Barrier: Deployable across smartphones, tablets, and low-cost PCs without friction.", level=1, bold=False, size=15, space_before=4)


    # ==========================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_common_header_footer(slide6, "RESEARCH AND REFERENCES", 6)

    tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.5))
    tf6 = tb6.text_frame
    tf6.word_wrap = True

    add_bullet(tf6, "• Details / Links of the reference and research work:", level=0, bold=True, size=17, space_before=0)
    add_bullet(tf6, "  - Self-Determination Theory (Deci & Ryan): Gamification drives intrinsic motivation via Autonomy (World exploration), Competence (XP/Leveling), and Relatedness (Social Leaderboards).", level=1, bold=False, size=14, space_before=5)
    add_bullet(tf6, "  - Bloom's Taxonomy Alignment: Dynamic Gemini AI quiz levels progress from basic recall to higher-order analytical reasoning.", level=1, bold=False, size=14, space_before=5)
    add_bullet(tf6, "  - Dual-Coding Theory (Paivio): Integrating ElevenLabs voice audio narration with Flame 2D visual game graphics dramatically improves information retention.", level=1, bold=False, size=14, space_before=5)
    add_bullet(tf6, "  - Flutter Framework Documentation: Multi-platform application architecture & UI engine (https://docs.flutter.dev/)", level=1, bold=False, size=14, space_before=5)
    add_bullet(tf6, "  - Flame 2D Game Engine Documentation: Open-source 2D engine for Flutter (https://flame-engine.org/)", level=1, bold=False, size=14, space_before=5)
    add_bullet(tf6, "  - Google Gemini API Documentation: Large language model & adaptive prompt engineering (https://ai.google.dev/)", level=1, bold=False, size=14, space_before=5)
    add_bullet(tf6, "  - ElevenLabs API Documentation: Conversational text-to-speech audio synthesis (https://elevenlabs.io/docs/)", level=1, bold=False, size=14, space_before=5)

    prs.save(output_pptx_path)
    print(f"[SUCCESS] Generated PPTX presentation at: {output_pptx_path}")

    # Export to PDF via PowerPoint COM
    try:
        import comtypes.client
        pdf_path = output_pptx_path.replace(".pptx", ".pdf")
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(output_pptx_path)
        presentation.SaveAs(pdf_path, 32) # 32 = ppSaveAsPDF
        presentation.Close()
        powerpoint.Quit()
        print(f"[SUCCESS] Exported PDF presentation at: {pdf_path}")
    except Exception as e:
        print(f"[NOTE] PDF export via COM skipped: {e}")

if __name__ == "__main__":
    out_file = os.path.abspath("HyperFusion_2026_M1racles_KnowledgeVerse.pptx")
    create_sih_presentation(out_file)
